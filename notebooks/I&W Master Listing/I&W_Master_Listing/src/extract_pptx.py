def extract_text_and_tables_from_pptx(pptx_path):
    from pptx import Presentation
    import re

    prs = Presentation(pptx_path)
    slides_content = []
    tables = []
    ip_like_data = []
    htoc_like_data = []
    htoc_pattern = re.compile(r'HTOC-\d{8}-\d{4}-[A-Z]')
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
                htoc_matches = htoc_pattern.findall(shape.text)
                htoc_like_data.extend(htoc_matches)
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                    for cell_text in row_data:
                        if (
                            "[" in cell_text and
                            "." in cell_text and
                            "]" in cell_text and
                            any(char.isdigit() for char in cell_text)
                        ):
                            ip_like_data.append(cell_text.strip())
                        htoc_matches = htoc_pattern.findall(cell_text)
                        htoc_like_data.extend(htoc_matches)
                tables.append(table_data)
        slides_content.append('\n'.join(slide_text))
    
    ip_like_data = list({item for item in ip_like_data if "[" in item and "]" in item and "." in item})
    htoc_like_data = list(set(htoc_like_data))
    
    return slides_content, tables, ip_like_data, htoc_like_data