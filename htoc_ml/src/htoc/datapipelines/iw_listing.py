"""Pull indicators out of finished I&W PDFs (or partner PPTX decks) into Excel."""
from __future__ import annotations

import argparse
import re
from datetime import datetime, timedelta
from pathlib import Path

import pandas as pd

from htoc.datapipelines.paths import env_path, share_root

HTOC_SERIAL = re.compile(r"HTOC-\d{8}-\d{4}-[A-Z]")
PARTNER_KEYWORDS = ("CDC", "NIH", "FDA", "HRSA", "VA", "CMS", "IHS", "DHA")
MASTER_COLUMNS = [
    "Partner",
    "I_W Description",
    "I_W#",
    "Indicator Disposition Code",
    "Indicator Disposition Code Description",
    "Secondary Indicator Disposition Code",
    "Secondary Indicator Disposition Code Description",
    "Tertiary Indicator Disposition Code",
    "Tertiary Indicator Disposition Code Description",
    "Comment",
    "Bi-Weekly Date",
    "Technique",
    "Malware",
    "Threat Actor",
    "Aliases",
    "Vulnerability",
    "Actor Tag",
    "Sector",
    "Real Organization",
    "Threat Actor Country",
    "I&W Serial",
    "Affected Partners",
]


def default_pdf_dir() -> Path:
    return env_path(
        "IW_LISTING_PDF_DIR",
        share_root() / "HTOC Reports" / "I&W Reports" / "3. Finished I&W Reports",
    )


def default_output_path() -> Path:
    return env_path(
        "IW_LISTING_OUTPUT",
        share_root() / "JA" / "IwListingTest" / "reported_iocs.xlsx",
    )


def default_pptx_root() -> Path:
    return env_path("IW_LISTING_PPTX_ROOT", share_root() / "Data_Analytics" / "I_W")


def default_master_xlsx() -> Path:
    return env_path(
        "IW_LISTING_MASTER_XLSX",
        share_root() / "JA" / "IwListingTest" / "Master.xlsx",
    )


def description_from_title(desc: str) -> str:
    match = re.search(r"[—-]\s*(.*?)(?:Seen|seen)", desc, re.DOTALL)
    if match:
        return match.group(1).strip()
    dash_idx = max(desc.rfind("—"), desc.rfind("-"))
    return desc[dash_idx + 1 :].strip() if dash_idx != -1 else ""


def serial_from_pdf_name(pdf_name: str) -> str:
    return Path(pdf_name).name.replace("HTOC-", "").replace(".pdf", "").replace(".PDF", "")


def recent_pdfs(folder: Path, lookback_weeks: int) -> list[Path]:
    cutoff = datetime.now() - timedelta(weeks=lookback_weeks)
    found: list[Path] = []
    if not folder.is_dir():
        return found
    for path in folder.iterdir():
        if path.is_file() and path.suffix.lower() == ".pdf":
            modified = datetime.fromtimestamp(path.stat().st_mtime)
            if modified >= cutoff:
                found.append(path)
    return sorted(found)


def extract_title_serial(first_page_text: str | None) -> str | None:
    if not first_page_text:
        return None
    text = first_page_text
    exec_idx = text.find("Executive Summary:")
    if exec_idx != -1:
        text = text[:exec_idx].strip()
    lines = [line.strip() for line in text.split("\n")]
    idx = next((i for i, line in enumerate(lines) if "Title/Serial Number:" in line), None)
    if idx is None:
        return None
    after = lines[idx].split("Title/Serial Number:", 1)[1].strip()
    collected = [after] if after else []
    for line in lines[idx + 1 :]:
        if line == "":
            break
        collected.append(line)
    return "\n".join(collected)


def normalize_indicator_table(frame: pd.DataFrame) -> pd.DataFrame:
    out = frame.copy()
    out.columns = [str(col).replace("\n", " ").strip() for col in out.columns]
    if "Indicator Type" in out.columns:
        out["Indicator Type"] = out["Indicator Type"].astype(str).str.replace("\n", " ", regex=False).str.strip()
    if "Indicator" in out.columns and "Type" in out.columns:
        out["Indicator Type"] = (out["Indicator"].fillna("") + " " + out["Type"].fillna("")).str.strip()
        out = out.drop(columns=["Indicator", "Type"])
    return out


def rows_from_indicator_table(frame: pd.DataFrame, serial: str) -> pd.DataFrame:
    cleaned = normalize_indicator_table(frame)
    indicator_cols = [
        col for col in cleaned.columns
        if "indicator" in str(col).lower() and " " not in str(col) and "\n" not in str(col)
    ]
    if not indicator_cols:
        return pd.DataFrame()
    type_col = "Indicator Type" if "Indicator Type" in cleaned.columns else None
    observed_col = "Observed By" if "Observed By" in cleaned.columns else None
    parts = []
    for col in indicator_cols:
        temp = cleaned[[col]].copy().rename(columns={col: "Indicator"})
        temp["Indicator"] = temp["Indicator"].astype(str).str.strip()
        temp["Type"] = cleaned[type_col] if type_col else pd.NA
        if observed_col:
            temp["Affected Partner(s)"] = cleaned[observed_col].astype(str).str.replace("\n", ", ")
        else:
            temp["Affected Partner(s)"] = pd.NA
        temp["I&W Serial"] = serial
        parts.append(temp)
    return pd.concat(parts, ignore_index=True) if parts else pd.DataFrame()


def combine_pdf_extractions(all_pdf_data: dict) -> pd.DataFrame:
    descriptions = {}
    frames: list[pd.DataFrame] = []
    for pdf_name, data in all_pdf_data.items():
        serial = serial_from_pdf_name(pdf_name)
        descriptions[serial] = description_from_title(str(data.get("title_serial_paragraph") or ""))
        for table_info in data.get("tables", []):
            frame = table_info.get("dataframe")
            if frame is None or frame.empty:
                continue
            rows = rows_from_indicator_table(frame, serial)
            if not rows.empty:
                frames.append(rows)
    if not frames:
        return pd.DataFrame()
    combined = pd.concat(frames, ignore_index=True)
    combined = combined[combined["Indicator"] != ""]
    combined["Description"] = combined["I&W Serial"].map(descriptions)
    return combined.reset_index(drop=True)


def extract_pdf(path: Path) -> dict:
    try:
        import pdfplumber
    except ImportError as exc:
        raise SystemExit("pdfplumber is not installed. py -m pip install pdfplumber") from exc
    pdf_data: dict = {"title_serial_paragraph": None, "tables": []}
    with pdfplumber.open(path) as pdf:
        if pdf.pages:
            pdf_data["title_serial_paragraph"] = extract_title_serial(pdf.pages[0].extract_text())
        for page_num, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables() or []
            for idx, table in enumerate(tables):
                if not table:
                    continue
                frame = pd.DataFrame(table[1:], columns=table[0])
                if "Indicator" in frame.columns and "Type" in frame.columns:
                    frame["Indicator"] = frame["Indicator"].replace("", pd.NA).ffill()
                    frame["Type"] = frame["Type"].replace("", pd.NA).ffill()
                    frame["Indicator Type"] = (
                        frame["Indicator"].astype(str) + " " + frame["Type"].astype(str)
                    ).str.strip()
                    frame = frame.drop(columns=["Indicator", "Type"])
                pdf_data["tables"].append({"page": page_num, "table_index": idx + 1, "dataframe": frame})
    return pdf_data


def write_tracker_excel(frame: pd.DataFrame, output_path: Path) -> Path:
    to_save = frame.drop(columns=["Description"], errors="ignore")
    cols = [c for c in to_save.columns if c != "Affected Partner(s)"]
    if "Affected Partner(s)" in to_save.columns:
        cols.append("Affected Partner(s)")
    to_save = to_save[cols].drop_duplicates()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        to_save.to_excel(writer, index=False)
        sheet = writer.sheets.get("Sheet1") or writer.book.active
        for column_cells in sheet.columns:
            max_length = max(len(str(cell.value)) if cell.value is not None else 0 for cell in column_cells)
            sheet.column_dimensions[column_cells[0].column_letter].width = max_length + 2
    return output_path


def is_defanged_ip(text: str) -> bool:
    return "[" in text and "]" in text and "." in text and any(ch.isdigit() for ch in text)


def extract_pptx_pairs(pptx_path: Path, partner: str | None) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("python-pptx is not installed. py -m pip install python-pptx") from exc
    prs = Presentation(str(pptx_path))
    pairs: list[dict] = []
    current_htoc = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                cells = [cell.text for cell in row.cells]
                htoc_matches = [cell for cell in cells if "HTOC-" in cell]
                if htoc_matches:
                    current_htoc = htoc_matches[0]
                ip_matches = [cell.strip() for cell in cells if is_defanged_ip(cell)]
                if current_htoc and ip_matches:
                    for ip in ip_matches:
                        pairs.append(
                            {"HTOC_Like_Data": current_htoc, "IP_Like_Data": ip, "Keyword": partner}
                        )
    return pairs


def newest_date_folder(root: Path) -> Path | None:
    folders = [p for p in root.iterdir() if p.is_dir() and p.name.isdigit() and len(p.name) == 8]
    return sorted(folders)[-1] if folders else None


def fill_master_sheet(file_path: Path, sheet_name: str, pairs: list[dict], date_str: str) -> Path:
    try:
        frame = pd.read_excel(file_path, sheet_name=sheet_name)
    except (FileNotFoundError, ValueError):
        frame = pd.DataFrame(columns=MASTER_COLUMNS)
    if frame.empty:
        frame = pd.DataFrame(columns=MASTER_COLUMNS)
    else:
        for col in MASTER_COLUMNS:
            if col not in frame.columns:
                frame[col] = ""
        frame = frame[MASTER_COLUMNS]
    added: set[tuple] = set()
    new_rows = []
    for pair in pairs:
        key = (pair["HTOC_Like_Data"], pair["IP_Like_Data"], pair["Keyword"])
        if key in added:
            continue
        added.add(key)
        new_rows.append(
            {
                "Partner": pair["Keyword"],
                "I_W Description": pair["HTOC_Like_Data"],
                "I_W#": pair["IP_Like_Data"],
                "Bi-Weekly Date": date_str,
                **{col: "" for col in MASTER_COLUMNS if col not in {"Partner", "I_W Description", "I_W#", "Bi-Weekly Date"}},
            }
        )
    if new_rows:
        frame = pd.concat([frame, pd.DataFrame(new_rows)], ignore_index=True)
    frame = frame.drop_duplicates(subset=["Partner", "I_W Description", "I_W#"], ignore_index=True)
    file_path.parent.mkdir(parents=True, exist_ok=True)
    writer_kw: dict = {"engine": "openpyxl", "mode": "a" if file_path.exists() else "w"}
    if writer_kw["mode"] == "a":
        writer_kw["if_sheet_exists"] = "replace"
    with pd.ExcelWriter(file_path, **writer_kw) as writer:
        frame.to_excel(writer, sheet_name=sheet_name, index=False)
    return file_path


def run_from_pdfs(*, pdf_dir: Path, output_path: Path, lookback_weeks: int) -> Path:
    paths = recent_pdfs(pdf_dir, lookback_weeks)
    if not paths:
        print(f"No PDFs in the last {lookback_weeks} week(s) under {pdf_dir}")
        return output_path
    all_pdf_data = {}
    for path in paths:
        print(f"Processing PDF: {path}")
        all_pdf_data[path.name] = extract_pdf(path)
    combined = combine_pdf_extractions(all_pdf_data)
    if combined.empty:
        print("No indicators found in stored tables.")
        return output_path
    write_tracker_excel(combined, output_path)
    print(f"Wrote {len(combined)} rows to {output_path}")
    return output_path


def run_from_pptx(*, pptx_root: Path, master_xlsx: Path, sheet_name: str) -> Path:
    folder = newest_date_folder(pptx_root) or pptx_root
    pptx_files = [
        p for p in folder.iterdir()
        if p.suffix.lower() == ".pptx" and any(k in p.name for k in PARTNER_KEYWORDS)
    ]
    pairs: list[dict] = []
    for path in pptx_files:
        partner = next((k for k in PARTNER_KEYWORDS if k in path.name), None)
        pairs.extend(extract_pptx_pairs(path, partner))
    date_str = datetime.strptime(folder.name, "%Y%m%d").strftime("%m/%d/%Y") if folder.name.isdigit() else datetime.now().strftime("%m/%d/%Y")
    fill_master_sheet(master_xlsx, sheet_name, pairs, date_str)
    print(f"Wrote {len(pairs)} pairs from {folder} to {master_xlsx}")
    return master_xlsx


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Extract indicators from finished I&W PDFs (default) or partner PPTX decks."
    )
    parser.add_argument("--from-pptx", action="store_true", help="Use the older partner-PPTX scan instead of PDFs.")
    parser.add_argument("--pdf-dir", default=str(default_pdf_dir()))
    parser.add_argument("--output", default=str(default_output_path()))
    parser.add_argument("--lookback-weeks", type=int, default=2)
    parser.add_argument("--pptx-root", default=str(default_pptx_root()))
    parser.add_argument("--master-xlsx", default=str(default_master_xlsx()))
    parser.add_argument("--sheet", default="Master Sheet")
    args = parser.parse_args(argv)
    if args.from_pptx:
        run_from_pptx(pptx_root=Path(args.pptx_root), master_xlsx=Path(args.master_xlsx), sheet_name=args.sheet)
    else:
        run_from_pdfs(pdf_dir=Path(args.pdf_dir), output_path=Path(args.output), lookback_weeks=args.lookback_weeks)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
